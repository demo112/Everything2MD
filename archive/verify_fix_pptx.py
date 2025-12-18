import sys
import os
from pathlib import Path
from src.core.converters.ppt import PptConverter
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)

def test_ppt_converter_fix():
    print("\n--- Testing PptConverter Fix ---")
    
    # Create dummy pptx if not exists
    input_path = Path("test_sample.pptx")
    if not input_path.exists():
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(input_path)
        print("Created test_sample.pptx")

    output_path = Path("output_verify/test_sample.md")
    
    converter = PptConverter()
    
    # Test 1: Check if _get_pptx2md_executable finds something
    exe_path = converter._get_pptx2md_executable()
    print(f"Found pptx2md executable at: {exe_path}")
    if "pptx2md" in exe_path and (os.path.exists(exe_path) or exe_path == "pptx2md"):
        print("Executable lookup seems valid.")
    else:
        print("WARNING: Executable lookup might be wrong.")

    # Test 2: Run conversion (This exercises the new str() conversion logic)
    try:
        result = converter.convert(input_path, output_path)
        print(f"Conversion successful: {result}")
    except Exception as e:
        print(f"Conversion failed: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_ppt_converter_fix()
