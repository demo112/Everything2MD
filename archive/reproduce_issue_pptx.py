import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def create_sample_pptx(path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle.text = "This is a test presentation."
    prs.save(path)
    print(f"Created sample PPTX at {path}")

def test_conversion_with_path_objects():
    print("\n--- Testing with Path objects (Expected to Fail) ---")
    try:
        from pptx2md.entry import convert as pptx_convert
        from pptx2md.types import ConversionConfig
    except ImportError:
        print("pptx2md not installed")
        return

    input_path = Path("test_sample.pptx")
    output_path = Path("output_path/test_sample.md")
    img_dir = Path("output_path/img")
    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    img_dir.mkdir(parents=True, exist_ok=True)

    config = ConversionConfig(
        pptx_path=input_path,
        output_path=output_path,
        image_dir=img_dir,
        title_path=None,
        image_width=None,
        disable_image=False,
        disable_wmf=False,
        disable_color=False,
        disable_escaping=False,
        disable_notes=False,
        enable_slides=False,
        try_multi_column=False,
        is_wiki=False,
        is_mdk=False,
        is_qmd=False,
        min_block_size=15,
        page=None,
        keep_similar_titles=False,
    )

    try:
        pptx_convert(config)
        print("SUCCESS: Converted with Path objects (Unexpected if bug exists)")
    except Exception as e:
        print(f"FAILURE: {type(e).__name__}: {e}")

def test_conversion_with_string_objects():
    print("\n--- Testing with String objects (Expected to Succeed) ---")
    try:
        from pptx2md.entry import convert as pptx_convert
        from pptx2md.types import ConversionConfig
    except ImportError:
        print("pptx2md not installed")
        return

    input_path = "test_sample.pptx"
    output_path = "output_str/test_sample.md"
    img_dir = "output_str/img"
    
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    Path(img_dir).mkdir(parents=True, exist_ok=True)

    config = ConversionConfig(
        pptx_path=input_path,
        output_path=output_path,
        image_dir=img_dir,
        title_path=None,
        image_width=None,
        disable_image=False,
        disable_wmf=False,
        disable_color=False,
        disable_escaping=False,
        disable_notes=False,
        enable_slides=False,
        try_multi_column=False,
        is_wiki=False,
        is_mdk=False,
        is_qmd=False,
        min_block_size=15,
        page=None,
        keep_similar_titles=False,
    )

    try:
        pptx_convert(config)
        print("SUCCESS: Converted with String objects")
    except Exception as e:
        print(f"FAILURE: {type(e).__name__}: {e}")

if __name__ == "__main__":
    create_sample_pptx("test_sample.pptx")
    test_conversion_with_path_objects()
    test_conversion_with_string_objects()
