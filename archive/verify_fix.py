import os
import sys
from pathlib import Path

# Force UTF-8 output
sys.stdout.reconfigure(encoding='utf-8')

# Add src to path
sys.path.append(os.path.abspath("src"))

try:
    from pptx import Presentation
    from core.converters.ppt import PptConverter as PPTConverter
    from core.utils import get_soffice_path
except ImportError as e:
    print(f"Import failed: {e}")
    sys.exit(1)

def create_test_pptx(path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Hello Verification"
    subtitle.text = "This is a test for PPTX conversion."
    
    prs.save(path)
    print(f"Created test PPTX at {path}")

def main():
    print("Starting verification...")
    
    # 1. Verify Dependencies
    try:
        import pptx2md
        print("✅ pptx2md is installed.")
    except ImportError:
        print("❌ pptx2md is NOT installed.")
        
    soffice = get_soffice_path()
    if soffice:
        print(f"✅ LibreOffice found at: {soffice}")
    else:
        print("⚠️ LibreOffice NOT found. PPT conversion will fail, but PPTX might work.")

    # 2. Test PPTX Conversion (Primary Path)
    test_pptx_path = Path("test_verify.pptx")
    output_md_path = Path("test_verify.md")
    
    create_test_pptx(test_pptx_path)
    
    converter = PPTConverter()
    
    try:
        print(f"Converting {test_pptx_path} to {output_md_path}...")
        converter.convert(test_pptx_path, output_md_path)
        
        if output_md_path.exists():
            content = output_md_path.read_text(encoding="utf-8")
            print(f"Output content length: {len(content)}")
            if "Hello Verification" in content:
                print("✅ PPTX Content verification PASSED!")
            else:
                print("❌ PPTX Content verification FAILED: Keyword not found.")
                print("Content preview:", content[:200])
        else:
            print("❌ Output file was not created.")
            
    except Exception as e:
        print(f"❌ Conversion failed with error: {e}")
        import traceback
        traceback.print_exc()

    # Clean up
    if test_pptx_path.exists():
        os.remove(test_pptx_path)
    if output_md_path.exists():
        os.remove(output_md_path)
    # Also clean up image dir if created
    img_dir = Path("test_verify") # pptx2md usually creates a dir with same name as output
    if img_dir.exists() and img_dir.is_dir():
        import shutil
        shutil.rmtree(img_dir)
        
if __name__ == "__main__":
    main()
