import pytest
import os
from pathlib import Path
from src.core.engine import ConversionEngine
from src.core.config import ConfigManager
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@pytest.mark.integration
def test_txt_conversion_flow(tmp_path):
    """
    Test the full flow from creating a file -> engine -> output
    Using .txt as it doesn't require external tools.
    """
    input_dir = tmp_path / "input"
    output_dir = tmp_path / "output"
    input_dir.mkdir()
    output_dir.mkdir()

    # 1. Prepare Data
    src_file = input_dir / "sample.txt"
    src_file.write_text("# Hello World\nThis is a test.")

    # 2. Setup Engine
    config = ConfigManager(str(tmp_path / "config.json"))
    engine = ConversionEngine(config)

    # 3. Run
    engine.run(str(input_dir), str(output_dir))

    # 4. Verify
    expected_out = output_dir / "sample.md"
    assert expected_out.exists()
    assert expected_out.read_text() == "# Hello World\nThis is a test."


@pytest.mark.integration
@pytest.mark.skipif(not HAS_DOCX, reason="python-docx not installed")
def test_docx_conversion_flow(tmp_path):
    """Test full flow for DOCX conversion"""
    input_dir = tmp_path / "input_docx"
    output_dir = tmp_path / "output_docx"
    input_dir.mkdir()
    output_dir.mkdir()

    # 1. Prepare Data
    src_file = input_dir / "test.docx"
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test paragraph.')
    doc.save(str(src_file))

    # 2. Setup Engine
    config = ConfigManager(str(tmp_path / "config.json"))
    # Ensure external tools are detected/configured if needed
    # (The Engine should auto-detect, or we assume default PATH)
    engine = ConversionEngine(config)

    # 3. Run
    engine.run(str(input_dir), str(output_dir))

    # 4. Verify
    expected_out = output_dir / "test.md"
    assert expected_out.exists()
    content = expected_out.read_text(encoding='utf-8')
    assert "Test Document" in content
    assert "This is a test paragraph" in content


@pytest.mark.integration
@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx not installed")
def test_pptx_conversion_flow(tmp_path):
    """Test full flow for PPTX conversion"""
    input_dir = tmp_path / "input_pptx"
    output_dir = tmp_path / "output_pptx"
    input_dir.mkdir()
    output_dir.mkdir()

    # 1. Prepare Data
    src_file = input_dir / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello PPTX"
    subtitle.text = "Subtitle Text"
    prs.save(str(src_file))

    # 2. Setup Engine
    config = ConfigManager(str(tmp_path / "config.json"))
    engine = ConversionEngine(config)

    # 3. Run
    engine.run(str(input_dir), str(output_dir))

    # 4. Verify
    expected_out = output_dir / "test.md"
    assert expected_out.exists()
    content = expected_out.read_text(encoding='utf-8')
    assert "Hello PPTX" in content
    assert "Subtitle Text" in content


@pytest.mark.integration
def test_gui_integration():
    """Verify GUI can import core modules (Smoke Test)"""
    try:
        from src.gui.main import Everything2MDGUI
        import tkinter as tk

        root = tk.Tk()
        app = Everything2MDGUI(root)
        root.destroy()
    except Exception as e:
        pytest.fail(f"GUI failed to initialize: {e}")
